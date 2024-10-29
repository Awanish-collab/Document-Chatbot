import streamlit as st
from pathlib import Path
import tempfile
import os
import win32com.client
import pythoncom
from pptx import Presentation
import time
from PIL import Image
import io
import base64

def convert_ppt_to_pptx(input_path):
    st.write("file path is: ", input_path)
    """Convert PPT file to PPTX"""
    pythoncom.CoInitialize()
    powerpoint = None
    temp_output_path = None
    
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(input_path, WithWindow=False)
        temp_output_path = str(Path(input_path).parent / f"temp_{int(time.time())}.pptx")
        presentation.SaveAs(temp_output_path, 24)  # 24 is the format for pptx
        presentation.Close()
        
        with open(temp_output_path, 'rb') as f:
            content = f.read()
        return content
    except Exception as e:
        st.error(f"Error converting file: {str(e)}")
        return None
    finally:
        if powerpoint:
            powerpoint.Quit()
        pythoncom.CoUninitialize()
        if temp_output_path and os.path.exists(temp_output_path):
            try:
                os.remove(temp_output_path)
            except:
                pass

def analyze_presentation(file_content):
    """Analyze presentation content including visual elements"""
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
        tmp_file.write(file_content)
        tmp_file_path = tmp_file.name

    try:
        prs = Presentation(tmp_file_path)
        total_slides = len(prs.slides)
        #st.write("Ppt details: ", prs)
        st.write("Total slides: ", total_slides)
        
        # Prepare HTML container with specific size and scroll
        html_content = """
        <div style="width: 100%; height: 400px; overflow-y: scroll; border: 1px solid black; padding: 10px;">
        """
        
        for i, slide in enumerate(prs.slides):
            # Convert slide to an image
            slide_image = slide_to_image(i, tmp_file_path)
            if slide_image:
                # Convert PIL image to base64 for HTML display
                buffered = io.BytesIO()
                slide_image.save(buffered, format="PNG")
                img_str = base64.b64encode(buffered.getvalue()).decode("utf-8")
                html_content += f'<img src="data:image/png;base64,{img_str}" style="width:100%; margin-bottom: 20px;" />'
            else:
                html_content += f'<p>Could not render slide {i + 1}</p>'

        html_content += "</div>"
        
        # Display scrollable container in Streamlit
        st.markdown(html_content, unsafe_allow_html=True)
        
    finally:
        try:
            os.unlink(tmp_file_path)
        except:
            pass

def slide_to_image(slide_number, ppt_path):
    """Convert a slide to an image using win32com.client"""
    pythoncom.CoInitialize()
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    presentation = powerpoint.Presentations.Open(ppt_path, WithWindow=False)
    
    try:
        slide = presentation.Slides(slide_number + 1)  # Adjusted indexing for PowerPoint
        # Specify output image path
        output_img_path = os.path.join(tempfile.gettempdir(), f"slide_{slide_number}.png")
        slide.Export(output_img_path, "PNG", 800, 600)
        
        with open(output_img_path, "rb") as f:
            img = Image.open(io.BytesIO(f.read()))
        return img
    except Exception as e:
        st.error(f"Error rendering slide {slide_number + 1}: {str(e)}")
        return None
    finally:
        presentation.Close()
        powerpoint.Quit()
        pythoncom.CoUninitialize()

def main():
    st.title("PowerPoint Content Analyzer")
    st.write("Upload a PowerPoint file (PPT or PPTX) to analyze its content.")

    uploaded_file = st.file_uploader(
        "Choose a PowerPoint file", 
        type=["ppt", "pptx"],
        help="Upload a PPT or PPTX file"
    )

    if uploaded_file is not None:
        file_extension = Path(uploaded_file.name).suffix.lower()
        
        # Create temporary directory for images
        with tempfile.TemporaryDirectory() as temp_dir:
            with st.spinner("Processing file..."):
                if file_extension == '.ppt':
                    with tempfile.NamedTemporaryFile(delete=False, suffix='.ppt') as tmp_file:
                        tmp_file.write(uploaded_file.getvalue())
                        tmp_file_path = tmp_file.name
                    
                    try:
                        pptx_content = convert_ppt_to_pptx(tmp_file_path)
                        if pptx_content is None:
                            st.error("Failed to convert PPT file.")
                            st.stop()
                    finally:
                        try:
                            os.unlink(tmp_file_path)
                        except:
                            pass
                else:
                    pptx_content = uploaded_file.getvalue()

                # Analyze the presentation content
                analyze_presentation(pptx_content)

if __name__ == "__main__":
    main()
