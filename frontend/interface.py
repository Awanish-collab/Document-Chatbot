'''import streamlit as st
import pdfplumber
from PIL import Image

def display_pdf(file_path, display_width=400):
    with pdfplumber.open(file_path) as pdf:
        num_pages = len(pdf.pages)
        st.write(f"Total Pages: {num_pages}")

        # Set up session state to keep track of the current page
        if 'page_num' not in st.session_state:
            st.session_state.page_num = 1

        # Input box for entering a page number directly
        page_input = st.number_input("Go to page:", min_value=1, max_value=num_pages, value=st.session_state.page_num)
        st.session_state.page_num = int(page_input)

        # Display the selected page with specified width
        page = pdf.pages[st.session_state.page_num - 1]
        page_image = page.to_image().original  # Convert to PIL image
        
        # Resize the image
        aspect_ratio = page_image.height / page_image.width
        resized_image = page_image.resize((display_width, int(display_width * aspect_ratio)))
        st.image(resized_image, use_column_width=False)

st.title("PDF Viewer")
pdf_file = st.file_uploader("Upload a PDF file", type="pdf")

if pdf_file:
    with open("uploaded_file.pdf", "wb") as f:
        f.write(pdf_file.getbuffer())
    display_pdf("uploaded_file.pdf", display_width=400)
'''

import streamlit as st
from pathlib import Path
import tempfile
import os
import win32com.client
import pythoncom
from pptx import Presentation
import io
from PIL import Image
import base64
import time
import numpy as np
from pptx.enum.shapes import MSO_SHAPE_TYPE
import matplotlib.pyplot as plt
from io import BytesIO

def extract_shape_image(shape, temp_dir):
    """Extract image from a shape"""
    try:
        if hasattr(shape, 'image'):
            # Get image data
            image_bytes = shape.image.blob
            # Save temporarily and convert to format streamlit can display
            img_path = os.path.join(temp_dir, f"shape_image_{time.time()}.png")
            with open(img_path, 'wb') as f:
                f.write(image_bytes)
            return img_path
    except Exception as e:
        st.error(f"Error extracting image: {str(e)}")
    return None

def extract_chart_image(shape, temp_dir):
    """Extract chart as image"""
    try:
        if shape.has_chart:
            # Get chart data
            chart = shape.chart
            chart_data = []
            
            # Extract basic chart data (this is simplified)
            if hasattr(chart, 'series'):
                for series in chart.series:
                    if hasattr(series, 'values'):
                        values = [v for v in series.values]
                        chart_data.append(values)

            # Create a simple matplotlib chart
            plt.figure(figsize=(10, 6))
            for data in chart_data:
                plt.plot(data)
            
            # Save the chart
            chart_path = os.path.join(temp_dir, f"chart_{time.time()}.png")
            plt.savefig(chart_path)
            plt.close()
            return chart_path
    except Exception as e:
        st.error(f"Error extracting chart: {str(e)}")
    return None

def convert_ppt_to_pptx(input_path):
    """Convert PPT file to PPTX"""
    pythoncom.CoInitialize()
    powerpoint = None
    temp_output_path = None
    
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(input_path)
        temp_output_path = str(Path(input_path).parent / f"temp_{int(time.time())}.pptx")
        presentation.SaveAs(temp_output_path, 24)
        presentation.Close()
        
        with open(temp_output_path, 'rb') as f:
            content = f.read()
        return content
    except Exception as e:
        st.error(f"Error converting file: {str(e)}")
        return None
    finally:
        if powerpoint:
            powerpoint.Quit()
        pythoncom.CoUninitialize()
        if temp_output_path and os.path.exists(temp_output_path):
            try:
                os.remove(temp_output_path)
            except:
                pass

def analyze_presentation(file_content, temp_dir):
    """Analyze presentation content including visual elements"""
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
        tmp_file.write(file_content)
        tmp_file_path = tmp_file.name

    try:
        prs = Presentation(tmp_file_path)
        content = []
        
        # Get presentation properties
        core_props = {
            "Author": prs.core_properties.author,
            "Created": prs.core_properties.created,
            "Modified": prs.core_properties.modified,
            "Title": prs.core_properties.title,
            "Subject": prs.core_properties.subject,
            "Keywords": prs.core_properties.keywords
        }
        
        total_slides = len(prs.slides)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = {
                "slide_number": slide_num,
                "layout": slide.slide_layout.name,
                "shapes": [],
                "texts": [],
                "tables": [],
                "charts": [],
                "images": []
            }
            
            for shape in slide.shapes:
                # Handle images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_path = extract_shape_image(shape, temp_dir)
                    if img_path:
                        slide_content["images"].append(img_path)
                
                # Handle charts
                elif shape.has_chart:
                    chart_path = extract_chart_image(shape, temp_dir)
                    if chart_path:
                        slide_content["charts"].append({
                            "path": chart_path,
                            "title": shape.chart.chart_title.text_frame.text if shape.chart.chart_title else "No Title"
                        })
                
                # Handle tables
                elif shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text)
                        table_data.append(row_data)
                    slide_content["tables"].append(table_data)
                
                # Handle text
                if shape.has_text_frame:
                    text_content = []
                    for paragraph in shape.text_frame.paragraphs:
                        para_info = {
                            "text": paragraph.text,
                            "level": paragraph.level,
                            "alignment": str(paragraph.alignment)
                        }
                        text_content.append(para_info)
                    slide_content["texts"].append(text_content)
            
            content.append(slide_content)
        
        return {
            "properties": core_props,
            "total_slides": total_slides,
            "slides": content
        }
        
    finally:
        try:
            os.unlink(tmp_file_path)
        except:
            pass

def main():
    st.title("PowerPoint Content Analyzer")
    st.write("Upload a PowerPoint file (PPT or PPTX) to analyze its content.")

    uploaded_file = st.file_uploader(
        "Choose a PowerPoint file", 
        type=["ppt", "pptx"],
        help="Upload a PPT or PPTX file"
    )

    if uploaded_file is not None:
        file_extension = Path(uploaded_file.name).suffix.lower()
        
        # Create temporary directory for images
        with tempfile.TemporaryDirectory() as temp_dir:
            with st.spinner("Processing file..."):
                if file_extension == '.ppt':
                    with tempfile.NamedTemporaryFile(delete=False, suffix='.ppt') as tmp_file:
                        tmp_file.write(uploaded_file.getvalue())
                        tmp_file_path = tmp_file.name
                    
                    try:
                        pptx_content = convert_ppt_to_pptx(tmp_file_path)
                        if pptx_content is None:
                            st.error("Failed to convert PPT file.")
                            st.stop()
                    finally:
                        try:
                            os.unlink(tmp_file_path)
                        except:
                            pass
                else:
                    pptx_content = uploaded_file.getvalue()

                # Analyze the presentation content
                content = analyze_presentation(pptx_content, temp_dir)

                if content:
                    # Display presentation properties
                    st.header("Presentation Properties")
                    props = content["properties"]
                    col1, col2 = st.columns(2)
                    with col1:
                        st.write("**Title:**", props["Title"] or "Not set")
                        st.write("**Author:**", props["Author"] or "Not set")
                        st.write("**Subject:**", props["Subject"] or "Not set")
                    with col2:
                        st.write("**Created:**", props["Created"])
                        st.write("**Modified:**", props["Modified"])
                        st.write("**Keywords:**", props["Keywords"] or "Not set")
                    
                    st.write(f"**Total Slides:** {content['total_slides']}")
                    
                    # Create tabs for slides
                    tabs = st.tabs([f"Slide {s['slide_number']}" for s in content['slides']])
                    
                    # Display content for each slide
                    for idx, slide in enumerate(content['slides']):
                        with tabs[idx]:
                            st.subheader(f"Slide {slide['slide_number']} - {slide['layout']}")
                            
                            # Create tabs for different content types
                            content_types = st.tabs(["Images", "Charts", "Text", "Tables"])
                            
                            # Images
                            with content_types[0]:
                                if slide["images"]:
                                    cols = st.columns(2)
                                    for i, img_path in enumerate(slide["images"]):
                                        with cols[i % 2]:
                                            try:
                                                image = Image.open(img_path)
                                                st.image(image, use_column_width=True)
                                            except Exception as e:
                                                st.error(f"Error displaying image: {str(e)}")
                                else:
                                    st.write("No images in this slide")
                            
                            # Charts
                            with content_types[1]:
                                if slide["charts"]:
                                    for chart in slide["charts"]:
                                        st.write(f"**{chart['title']}**")
                                        try:
                                            image = Image.open(chart["path"])
                                            st.image(image, use_column_width=True)
                                        except Exception as e:
                                            st.error(f"Error displaying chart: {str(e)}")
                                else:
                                    st.write("No charts in this slide")
                            
                            # Text Content
                            with content_types[2]:
                                if slide["texts"]:
                                    for text_block in slide["texts"]:
                                        for paragraph in text_block:
                                            st.markdown(f"- {paragraph['text']}")
                                            st.caption(f"Level: {paragraph['level']}, Alignment: {paragraph['alignment']}")
                                else:
                                    st.write("No text content in this slide")
                            
                            # Tables
                            with content_types[3]:
                                if slide["tables"]:
                                    for idx, table in enumerate(slide["tables"], 1):
                                        st.write(f"Table {idx}:")
                                        st.table(table)
                                else:
                                    st.write("No tables in this slide")

if __name__ == "__main__":
    main()
