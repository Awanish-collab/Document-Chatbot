import streamlit as st
from PIL import Image
import os
import time
from PyPDF2 import PdfReader
from docx import Document
import speech_recognition as sr
from moviepy.editor import VideoFileClip
from pptx import Presentation
import tempfile
from collections import defaultdict
import hashlib
from datetime import datetime
import base64

class FileStatus:
    PROCESSING = "processing"
    COMPLETED = "completed"
    FAILED = "failed"

class DocumentProcessor:
    @staticmethod
    def get_file_hash(file_content):
        return hashlib.md5(file_content).hexdigest()

    @staticmethod
    def extract_text_from_pdf(file):
        try:
            pdf_reader = PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            if not text.strip():
                raise ValueError("No text content extracted from PDF")
            return text
        except Exception as e:
            raise Exception(f"PDF processing error: {str(e)}")

    @staticmethod
    def extract_text_from_docx(file):
        try:
            doc = Document(file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            if not text.strip():
                raise ValueError("No text content extracted from DOCX")
            return text
        except Exception as e:
            raise Exception(f"DOCX processing error: {str(e)}")

    @staticmethod
    def extract_text_from_pptx(file):
        try:
            ppt = Presentation(file)
            text = ""
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            if not text.strip():
                raise ValueError("No text content extracted from PPTX")
            return text
        except Exception as e:
            raise Exception(f"PPTX processing error: {str(e)}")

    @staticmethod
    def extract_audio_from_video(video_path):
        try:
            with VideoFileClip(video_path) as video:
                audio = video.audio
                if audio is None:
                    raise ValueError("No audio found in video file")
                return audio
        except Exception as e:
            raise Exception(f"Video processing error: {str(e)}")

    @staticmethod
    def extract_text_from_audio(audio_file):
        try:
            recognizer = sr.Recognizer()
            with sr.AudioFile(audio_file) as source:
                audio = recognizer.record(source)
            text = recognizer.recognize_google(audio)
            if not text.strip():
                raise ValueError("No text content extracted from audio")
            return text
        except Exception as e:
            raise Exception(f"Audio processing error: {str(e)}")

def initialize_session_state():
    if 'messages' not in st.session_state:
        st.session_state.messages = []
    if 'uploaded_files' not in st.session_state:
        st.session_state.uploaded_files = defaultdict(list)
    if 'file_contents' not in st.session_state:
        st.session_state.file_contents = {}
    if 'file_hashes' not in st.session_state:
        st.session_state.file_hashes = {}
    if 'file_status' not in st.session_state:
        st.session_state.file_status = {}
    if 'file_metadata' not in st.session_state:
        st.session_state.file_metadata = {}
    if 'file_data' not in st.session_state:
        st.session_state.file_data = {}

def truncate_filename(filename, max_length=15):
    name, ext = os.path.splitext(filename)
    if len(name) > max_length:
        return name[:max_length] + "..." + ext
    return filename

def get_file_type_icon(file_type):
    icons = {
        'pdf': '📄',
        'docx': '📝',
        'mp4': '🎥',
        'pptx': '📊',
        'mp3': '🎵',
        'wav': '🎵'
    }
    return icons.get(file_type, '📎')

def get_file_type(filename):
    return filename.split('.')[-1].lower()

def format_size(size):
    for unit in ['B', 'KB', 'MB', 'GB']:
        if size < 1024:
            return f"{size:.2f} {unit}"
        size /= 1024
    return f"{size:.2f} TB"

def get_file_details_html(filename):
    metadata = st.session_state.file_metadata.get(filename, {})
    status = st.session_state.file_status.get(filename)
    
    if status == FileStatus.COMPLETED:
        details = f"""
        <div style='background-color: white; padding: 10px; border-radius: 5px; box-shadow: 0 2px 4px rgba(0,0,0,0.1);'>
            <p><strong>✅ Status:</strong> Completed</p>
            <p><strong>📅 Processed At:</strong> {metadata.get('processed_at', 'N/A')}</p>
            <p><strong>📊 File Size:</strong> {format_size(metadata.get('size', 0))}</p>
            <p><strong>📝 Content Length:</strong> {metadata.get('content_length', 0):,} characters</p>
            <p><strong>🏷️ File Type:</strong> {metadata.get('type', 'unknown').upper()}</p>
        </div>
        """
    else:
        details = f"""
        <div style='background-color: white; padding: 10px; border-radius: 5px; box-shadow: 0 2px 4px rgba(0,0,0,0.1);'>
            <p><strong>❌ Status:</strong> Failed</p>
            <p><strong>⚠️ Error:</strong> {metadata.get('error_message', 'Unknown error')}</p>
            <p><strong>📅 Failed At:</strong> {metadata.get('failed_at', 'N/A')}</p>
        </div>
        """
    
    return details

def process_file(file):
    file_type = get_file_type(file.name)
    st.session_state.file_status[file.name] = FileStatus.PROCESSING
    
    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            if file_type in ['mp4', 'mp3', 'wav']:
                tmp_path = os.path.join(temp_dir, file.name)
                with open(tmp_path, 'wb') as tmp_file:
                    tmp_file.write(file.getvalue())

                if file_type == 'mp4':
                    audio = DocumentProcessor.extract_audio_from_video(tmp_path)
                    audio_path = os.path.join(temp_dir, 'audio.wav')
                    audio.write_audiofile(audio_path)
                    content = DocumentProcessor.extract_text_from_audio(audio_path)
                else:
                    content = DocumentProcessor.extract_text_from_audio(tmp_path)
            else:
                if file_type == 'pdf':
                    content = DocumentProcessor.extract_text_from_pdf(file)
                elif file_type == 'docx':
                    content = DocumentProcessor.extract_text_from_docx(file)
                elif file_type == 'pptx':
                    content = DocumentProcessor.extract_text_from_pptx(file)
                else:
                    raise ValueError(f"Unsupported file type: {file_type}")

        if not content or not content.strip():
            raise ValueError("No content extracted from file")

        st.session_state.file_status[file.name] = FileStatus.COMPLETED
        st.session_state.file_metadata[file.name] = {
            'processed_at': datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            'size': len(file.getvalue()),
            'content_length': len(content),
            'type': file_type
        }
        
        # Store the file data for preview
        st.session_state.file_data[file.name] = file.getvalue()
        
        return content

    except Exception as e:
        st.session_state.file_status[file.name] = FileStatus.FAILED
        st.session_state.file_metadata[file.name] = {
            'error_message': str(e),
            'failed_at': datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            'type': file_type
        }
        raise

def create_file_preview(filename):
    file_type = get_file_type(filename)
    file_data = st.session_state.file_data.get(filename)
    
    if not file_data:
        return "File preview not available"
    
    if file_type in ['mp4', 'mp3', 'wav']:
        # Create a data URL for audio/video files
        mime_types = {
            'mp4': 'video/mp4',
            'mp3': 'audio/mpeg',
            'wav': 'audio/wav'
        }
        data_url = f"data:{mime_types[file_type]};base64,{base64.b64encode(file_data).decode()}"
        if file_type == 'mp4':
            return f'<video controls width="100%"><source src="{data_url}" type="video/mp4"></video>'
        else:
            return f'<audio controls><source src="{data_url}" type="{mime_types[file_type]}"></audio>'
    
    elif file_type == 'pdf':
        data_url = f"data:application/pdf;base64,{base64.b64encode(file_data).decode()}"
        #return f'<embed src="{data_url}" width="100%" height="500px" type="">'
        return f'<iframe src="{data_url}" height="500" width="100%" title="application/pdf"></iframe>'
    
    else:
        # For other file types, show the extracted text content
        content = st.session_state.file_contents.get(filename, "Content preview not available")
        return f'<pre style="max-height: 300px; overflow-y: auto;">{content[:1000]}...</pre>'

def handle_chat_input(prompt):
    if not st.session_state.file_contents:
        return "Please upload and process some documents before asking questions."

    # Clear previous messages when new files are uploaded
    st.session_state.messages = []
    
    response = "Based on the processed documents:\n\n"
    for file_type, filenames in st.session_state.uploaded_files.items():
        successfully_processed = [
            f for f in filenames 
            if st.session_state.file_status.get(f) == FileStatus.COMPLETED
        ]
        if successfully_processed:
            response += f"\n{get_file_type_icon(file_type)} {file_type.upper()} files:\n"
            for filename in successfully_processed:
                content_preview = st.session_state.file_contents[filename][:100] + "..."
                response += f"- {truncate_filename(filename)}\n  Preview: {content_preview}\n"
    
    return response

def main():
    st.set_page_config(
        page_title="Document Q&A Chatbot",
        layout="wide",
        initial_sidebar_state="expanded"
    )

    # Custom CSS for hover effect
    st.markdown("""
        <style>
        .file-container {
            position: relative;
            padding: 10px;
            margin: 5px 0;
            border-radius: 5px;
            background-color: #f0f2f6;
        }
        .file-container:hover {
            background-color: #e6e9ef;
        }
        .hover-details {
            display: none;
            position: absolute;
            right: 100%;
            top: 0;
            z-index: 1000;
            width: 300px;
            margin-right: 10px;
        }
        .file-container:hover .hover-details {
            display: block;
        }
        .preview-button {
            background-color: #4CAF50;
            color: white;
            padding: 5px 10px;
            border: none;
            border-radius: 3px;
            cursor: pointer;
        }
        .preview-button:hover {
            background-color: #45a049;
        }
        </style>
    """, unsafe_allow_html=True)

    initialize_session_state()

    with st.sidebar:
        st.title("📁 Document Manager")
        
        uploaded_files = st.file_uploader(
            "Upload your documents",
            type=['pdf', 'docx', 'mp4', 'pptx', 'mp3', 'wav'],
            accept_multiple_files=True
        )

        if uploaded_files:
            for file in uploaded_files:
                file_content = file.read()
                file_hash = DocumentProcessor.get_file_hash(file_content)
                file.seek(0)
                
                if file.name in st.session_state.file_hashes and st.session_state.file_hashes[file.name] == file_hash:
                    st.warning(f'{truncate_filename(file.name)} is a duplicate and was skipped.')
                    continue
                
                st.session_state.file_hashes[file.name] = file_hash
                file_type = get_file_type(file.name)
                
                with st.spinner(f'Processing {truncate_filename(file.name)}...'):
                    try:
                        text_content = process_file(file)
                        if file_type not in st.session_state.uploaded_files:
                            st.session_state.uploaded_files[file_type] = []
                        if file.name not in st.session_state.uploaded_files[file_type]:
                            st.session_state.uploaded_files[file_type].append(file.name)
                        st.session_state.file_contents[file.name] = text_content
                        st.success(f'Successfully processed {truncate_filename(file.name)}')
                    except Exception as e:
                        st.error(f'Error processing {truncate_filename(file.name)}: {str(e)}')
                        if file.name in st.session_state.file_hashes:
                            del st.session_state.file_hashes[file.name]

        if st.session_state.uploaded_files:
            st.write("### Processed Files")
            
            file_types = list(st.session_state.uploaded_files.keys())
            if file_types:
                tabs = st.tabs([f"{get_file_type_icon(ft)} {ft.upper()}" for ft in file_types])
                
                for tab, file_type in zip(tabs, file_types):
                    with tab:
                        for filename in st.session_state.uploaded_files[file_type]:
                            # Create a container for each file with hover effect
                            st.markdown(f"""
                                <div class="file-container">
                                    <div class="hover-details">
                                        {get_file_details_html(filename)}
                                    </div>
                                    <div style="display: flex; justify-content: space-between; align-items: center;">
                                        <span>{get_file_type_icon(file_type)} {truncate_filename(filename)}</span>
                                        <button class="preview-button" onclick="showPreview('{filename}')">
                                            👁️ Preview
                                        </button>
                                    </div>
                                </div>
                            """, unsafe_allow_html=True)
                            
                            # Create a unique key for each file's preview expander
                            preview_key = f"preview_{filename}"
                            if st.button("👁️ Preview", key=preview_key):
                                st.markdown("### File Preview")
                                preview_html = create_file_preview(filename)
                                st.markdown(preview_html, unsafe_allow_html=True)

            st.write("### Summary")
            total_files = sum(len(files) for files in st.session_state.uploaded_files.values())
            completed_files = sum(1 for status in st.session_state.file_status.values() 
                                if status == FileStatus.COMPLETED)
            failed_files = sum(1 for status in st.session_state.file_status.values() 
                             if status == FileStatus.FAILED)
            
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("Total Files", total_files)
            with col2:
                st.metric("Successfully Processed", completed_files)
            with col3:
                st.metric("Failed", failed_files)

    st.title("💬 Document Q&A Chatbot")

    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.write(message["content"])

    if prompt := st.chat_input("Ask a question about your documents"):
        with st.chat_message("user"):
            st.write(prompt)

        with st.chat_message("assistant"):
            message_placeholder = st.empty()
            full_response = handle_chat_input(prompt)
            
            displayed_response = ""
            for chunk in full_response.split():
                displayed_response += chunk + " "
                time.sleep(0.05)
                message_placeholder.markdown(displayed_response + "▌")
            message_placeholder.markdown(displayed_response)
        
        st.session_state.messages.append({
            "role": "assistant", 
            "content": full_response
        })

if __name__ == "__main__":
    main()